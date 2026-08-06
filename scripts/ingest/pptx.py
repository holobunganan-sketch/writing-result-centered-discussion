from __future__ import annotations

import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from .base import ExtractedUnit, unit


def _fallback(path: Path) -> list[ExtractedUnit]:
    units: list[ExtractedUnit] = []
    with zipfile.ZipFile(path) as archive:
        slides = sorted((n for n in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)), key=lambda n: int(re.search(r"\d+", n).group()))
        for idx, name in enumerate(slides, 1):
            root = ET.fromstring(archive.read(name))
            text = " ".join(x.text or "" for x in root.iter() if x.tag.rsplit("}", 1)[-1] == "t")
            if found := unit(f"slide {idx} shape 1", text, format="pptx"):
                units.append(found)
    return units


def extract_pptx(path: Path) -> list[ExtractedUnit]:
    try:
        from pptx import Presentation  # type: ignore
        presentation = Presentation(path)
        units: list[ExtractedUnit] = []
        for s_idx, slide in enumerate(presentation.slides, 1):
            for sh_idx, shape in enumerate(slide.shapes, 1):
                if getattr(shape, "has_text_frame", False):
                    text = shape.text
                    if found := unit(f"slide {s_idx} shape {sh_idx}", text, format="pptx", shape_name=shape.name):
                        units.append(found)
                if getattr(shape, "has_table", False):
                    for r_idx, row in enumerate(shape.table.rows, 1):
                        for c_idx, cell in enumerate(row.cells, 1):
                            if found := unit(f"slide {s_idx} table {sh_idx} row {r_idx} col {c_idx}", cell.text, format="pptx-table"):
                                units.append(found)
            try:
                notes = slide.notes_slide.notes_text_frame.text
                if found := unit(f"slide {s_idx} notes", notes, format="pptx-notes"):
                    units.append(found)
            except Exception:
                pass
        return units or _fallback(path)
    except Exception:
        return _fallback(path)
